import json
import os
import shutil
from collections import Counter

import cv2
import numpy as np
from PIL import ImageFont, Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .utils import extract_background_color, extract_font_color, fill_bbox_with_bg, get_projection_segments


class Character:
    """A simple data class to hold information about a detected character."""

    def __init__(self, bbox, color, line_index):
        self.bbox = bbox
        self.color = color
        self.line_index = line_index
        self.font_size = 0
        self.bold = False
        self.text = ""

    def __repr__(self):
        return (f"Character(text='{self.text}', bbox={self.bbox}, color={self.color}, line={self.line_index}, "
                f"size={self.font_size}, bold={self.bold})")


class PageContext:
    def __init__(self, page_image, coords, slide):
        self.slide = slide
        self.original_image = page_image.copy()
        self.background_image = page_image.copy()
        self.coords = coords
        self.elements = []
        self.raw_chars = []
        self.corrected_chars = []

    def add_element_bbox_for_cleanup(self, bbox):
        """Register a bounding box to be inpainted on the background image."""
        if bbox:
            px_box = [int(v * (self.coords['img_w'] / self.coords['json_w'] if i % 2 == 0 else self.coords['img_h'] / self.coords['json_h'])) for i, v in enumerate(bbox)]
            fill_bbox_with_bg(self.background_image, px_box)

    def add_processed_element(self, elem_type, data):
        """Store a fully processed element ready for rendering."""
        self.elements.append({'type': elem_type, 'data': data})

    def add_characters(self, raw_chars, corrected_chars):
        if raw_chars: self.raw_chars.extend(raw_chars)
        if corrected_chars: self.corrected_chars.extend(corrected_chars)

    def generate_debug_images(self, page_index, generator_instance):
        """Generate and save debug images for the page."""
        generator_instance._draw_debug_boxes_for_page(self.original_image, self.raw_chars, self.coords, f"tmp/page_{page_index}_raw.png")
        generator_instance._draw_debug_boxes_for_page(self.original_image, self.corrected_chars, self.coords, f"tmp/page_{page_index}_corrected.png")

    def render_to_slide(self, generator_instance):
        """Render all processed elements onto the PowerPoint slide."""
        # 1. Render the cleaned background
        bg_path = f"temp_bg_{id(self.slide)}.png"
        cv2.imwrite(bg_path, cv2.cvtColor(self.background_image, cv2.COLOR_RGB2BGR))
        w_pts, h_pts = generator_instance.prs.slide_width, generator_instance.prs.slide_height
        self.slide.shapes.add_picture(bg_path, Pt(0), Pt(0), w_pts, h_pts)
        os.remove(bg_path)

        # 2. Render all image elements first
        for elem in self.elements:
            if elem['type'] == 'image':
                generator_instance._add_picture_from_bbox(self.slide, elem['data']['bbox'], self.original_image, self.coords, elem['data']['text_elements'])

        # 3. Render all text elements on top
        for elem in self.elements:
            if elem['type'] == 'text':
                generator_instance._render_text_from_data(self.slide, elem['data'])


class PPTGenerator:
    def __init__(self, output_path, remove_watermark=True):
        self.prs = Presentation()
        self.output_path = output_path
        self.remove_watermark = remove_watermark
        self.debug_images = False # Will be set in process_page
        for i in range(len(self.prs.slides) - 1, -1, -1):
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]

    def cap_size(self, w_pts, h_pts):
        MAX_PTS = 56 * 72
        if w_pts > MAX_PTS or h_pts > MAX_PTS:
            scale = MAX_PTS / max(w_pts, h_pts)
            w_pts, h_pts = w_pts * scale, h_pts * scale
        return w_pts, h_pts

    def set_slide_size(self, width_px, height_px, dpi=72):
        w_pts, h_pts = self.cap_size(width_px * 72 / dpi, height_px * 72 / dpi)
        self.prs.slide_width, self.prs.slide_height = Pt(w_pts), Pt(h_pts)

    def add_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _get_bbox_intersection(self, bbox1, bbox2):
        x1, y1 = max(bbox1[0], bbox2[0]), max(bbox1[1], bbox2[1])
        x2, y2 = min(bbox1[2], bbox2[2]), min(bbox1[3], bbox2[3])
        return [x1, y1, x2, y2] if x1 < x2 and y1 < y2 else None

    def _create_textbox(self, slide, bbox, coords):
        x1, y1, x2, y2 = bbox
        return slide.shapes.add_textbox(
            Pt(x1 * coords['scale_x']), Pt(y1 * coords['scale_y']),
            Pt((x2 - x1) * coords['scale_x']), Pt((y2 - y1) * coords['scale_y'])
        )

    def _get_line_ranges(self, page_image, bbox, coords):
        x1, y1, x2, y2 = bbox
        px1, py1, px2, py2 = int(x1 * coords['img_w'] / coords['json_w']), int(
            y1 * coords['img_h'] / coords['json_h']), int(x2 * coords['img_w'] / coords['json_w']), int(
            y2 * coords['img_h'] / coords['json_h'])
        h, w = page_image.shape[:2];
        px1, py1, px2, py2 = max(0, px1), max(0, py1), min(w, px2), min(h, py2)
        if px2 <= px1 or py2 <= py1: return []
        roi = page_image[py1:py2, px1:px2]
        bg_color = extract_background_color(page_image, [px1, py1, px2, py2])
        font_color, _, _ = extract_font_color(page_image, [px1, py1, px2, py2], bg_color)
        initial_lines = get_projection_segments(roi, font_color, axis=1)
        line_infos = []
        scale_y = (y2 - y1) / roi.shape[0] if roi.shape[0] > 0 else 0
        for start_y, end_y in initial_lines:
            line_pixel_bbox = [px1, py1 + start_y, px2, py1 + end_y]
            line_bg = extract_background_color(page_image, line_pixel_bbox)
            line_fg, _, _ = extract_font_color(page_image, line_pixel_bbox, line_bg)
            line_infos.append({'range': [y1 + start_y * scale_y, y1 + end_y * scale_y], 'color': line_fg,
                               'pixel_range': (start_y, end_y)})
        if not line_infos: return []
        avg_line_height = np.mean([info['pixel_range'][1] - info['pixel_range'][0] for info in line_infos])
        recovered_lines = []
        sorted_lines = sorted(line_infos, key=lambda x: x['pixel_range'][0])
        all_gaps = [(0, sorted_lines[0]['pixel_range'][0])] + [
            (sorted_lines[i]['pixel_range'][1], sorted_lines[i + 1]['pixel_range'][0]) for i in
            range(len(sorted_lines) - 1)] + [(sorted_lines[-1]['pixel_range'][1], roi.shape[0])]
        for gap_start, gap_end in all_gaps:
            if (gap_end - gap_start) > avg_line_height * 0.8:
                gap_bbox = [px1, py1 + gap_start, px2, py1 + gap_end]
                gap_bg = extract_background_color(page_image, gap_bbox)
                new_font_color, x_prop, y_prop = extract_font_color(page_image, gap_bbox, gap_bg)
                if y_prop > x_prop * 1.2 and np.linalg.norm(np.array(new_font_color) - np.array(font_color)) > 50:
                    gap_roi = roi[gap_start:gap_end, :];
                    gap_pixels = gap_roi.reshape(-1, 3)
                    gap_diff = np.linalg.norm(gap_pixels - new_font_color, axis=1)
                    gap_mask = (gap_diff < 40).reshape(gap_roi.shape[:2])
                    gap_row_counts = np.sum(gap_mask, axis=1)
                    in_gap_line, gap_line_start = False, 0
                    for y, count in enumerate(gap_row_counts):
                        if count > 1 and not in_gap_line:
                            in_gap_line, gap_line_start = True, y
                        elif count < 1 and in_gap_line:
                            in_gap_line = False
                            if y - gap_line_start > 3:
                                abs_start, abs_end = gap_start + gap_line_start, gap_start + y
                                recovered_lines.append({'range': [y1 + abs_start * scale_y, y1 + abs_end * scale_y],
                                                        'color': new_font_color, 'pixel_range': (abs_start, abs_end)})
                    if in_gap_line:
                        abs_start, abs_end = gap_start + gap_line_start, gap_start + len(gap_row_counts)
                        recovered_lines.append(
                            {'range': [y1 + abs_start * scale_y, y1 + abs_end * scale_y], 'color': new_font_color,
                             'pixel_range': (abs_start, abs_end)})
        if recovered_lines:
            line_infos.extend(recovered_lines)
            line_infos.sort(key=lambda x: x['range'][0])

        if len(line_infos) > 1:
            avg_line_height = np.mean([info['pixel_range'][1] - info['pixel_range'][0] for info in line_infos])
            merged_lines = [line_infos[0]]
            for i in range(1, len(line_infos)):
                prev_line = merged_lines[-1]
                curr_line = line_infos[i]
                gap = curr_line['pixel_range'][0] - prev_line['pixel_range'][1]
                if gap >= 0 and gap <= max(avg_line_height * 0.05, 1):
                    prev_height = prev_line['pixel_range'][1] - prev_line['pixel_range'][0]
                    curr_height = curr_line['pixel_range'][1] - curr_line['pixel_range'][0]
                    new_pixel_range = (prev_line['pixel_range'][0], curr_line['pixel_range'][1])
                    new_range = [prev_line['range'][0], curr_line['range'][1]]
                    new_color = curr_line['color'] if curr_height > prev_height else prev_line['color']
                    merged_lines[-1] = {'range': new_range, 'color': new_color, 'pixel_range': new_pixel_range}
                else:
                    merged_lines.append(curr_line)
            line_infos = merged_lines

        for i, info in enumerate(line_infos):
            # Define the top of the search area as the bottom of the previous line.
            info['search_top_y'] = line_infos[i - 1]['range'][1] if i > 0 else bbox[1]

        return line_infos

    def _detect_raw_characters(self, page_image, line_infos, bbox, coords):
        char_objects = []
        for i, info in enumerate(line_infos):
            tight_bbox = [bbox[0], info['range'][0], bbox[2], info['range'][1]]
            search_top_y = info['search_top_y']
            char_objects.extend(
                self._detect_characters_from_line(page_image, tight_bbox, search_top_y, coords, info['color'], i))
        return char_objects

    def _detect_characters_from_line(self, page_image, tight_bbox, search_top_y, coords, line_color, line_index):
        x1, y1, x2, y2 = tight_bbox
        # Convert JSON coordinates to pixel coordinates for the tight box and the search boundary
        px1 = int(x1 * coords['img_w'] / coords['json_w'])
        py1 = int(y1 * coords['img_h'] / coords['json_h'])
        px2 = int(x2 * coords['img_w'] / coords['json_w'])
        py2 = int(y2 * coords['img_h'] / coords['json_h'])
        search_top_py = int(search_top_y * coords['img_h'] / coords['json_h'])

        h, w = page_image.shape[:2]
        px1, py1, px2, py2 = max(0, px1), max(0, py1), min(w, px2), min(h, py2)
        search_top_py = max(0, search_top_py)

        if px2 <= px1 or py2 <= py1:
            return []

        # Define the single, consistent scaling factors for this line based on the tight box.
        scale_x = (x2 - x1) / (px2 - px1) if (px2 - px1) > 0 else 0
        scale_y = (y2 - y1) / (py2 - py1) if (py2 - py1) > 0 else 0

        # Define the region of interest strictly for the tight line box for primary character segmentation.
        tight_roi = page_image[py1:py2, px1:px2]

        # Segment main characters within the TIGHT ROI.
        all_chars = self._segment_characters_in_roi(
            tight_roi, tight_bbox, line_color, line_index, scale_x
        )

        if not all_chars:
            return []
        sorted_chars = sorted(all_chars, key=lambda c: c.bbox[0])

        # Find gaps between characters to search for text of a different color.
        gaps, last_x2 = [], tight_bbox[0]
        for char in sorted_chars:
            if char.bbox[0] > last_x2: gaps.append((last_x2, char.bbox[0]))
            last_x2 = char.bbox[2]
        if tight_bbox[2] > last_x2: gaps.append((last_x2, tight_bbox[2]))

        recovered_chars = []
        scale_x_inv = (px2 - px1) / (x2 - x1) if (x2 - x1) > 0 else 0
        for gap_x1, gap_x2 in gaps:
            gap_px1 = px1 + int((gap_x1 - x1) * scale_x_inv)
            gap_px2 = px1 + int((gap_x2 - x1) * scale_x_inv)

            # Shrink the search box by 5 pixels on each side to avoid edge artifacts from the primary font.
            gap_px1 += 5
            gap_px2 -= 5

            if gap_px2 - gap_px1 < 30: continue
            gap_roi = page_image[search_top_py:py2, gap_px1:gap_px2]
            if gap_roi.size == 0: continue

            gap_bg = extract_background_color(page_image, [gap_px1, search_top_py, gap_px2, py2])
            new_font_color, x_prop, y_prop = extract_font_color(page_image,
                                                                [gap_px1, search_top_py, gap_px2, py2], gap_bg)

            if max(x_prop, y_prop) > 0.15 and np.linalg.norm(np.array(new_font_color) - np.array(line_color)) > 50:
                segments = get_projection_segments(gap_roi, new_font_color, axis=1)

                if segments:
                    # Find the tallest segment, as it's the most likely candidate for the actual line of text.
                    best_segment = max(segments, key=lambda s: s[1] - s[0])
                    segment_height = best_segment[1] - best_segment[0]

                    if segment_height >= 8:
                        local_py1 = best_segment[0]
                        adjusted_roi_py1 = search_top_py + local_py1
                        adjusted_gap_roi = page_image[adjusted_roi_py1:py2, gap_px1:gap_px2]

                        if adjusted_gap_roi.size > 0:
                            new_tight_y1 = search_top_y + (adjusted_roi_py1 - search_top_py) * scale_y
                            new_tight_bbox = [gap_x1, new_tight_y1, gap_x2, y2]
                            recovered_chars.extend(
                                self._segment_characters_in_roi(
                                    adjusted_gap_roi, new_tight_bbox, new_font_color, line_index, scale_x
                                )
                            )
        if recovered_chars:
            all_chars.extend(recovered_chars)
            all_chars.sort(key=lambda c: c.bbox[0])
        return all_chars

    def _segment_characters_in_roi(self, roi, tight_bbox, color, line_index, scale_x):
        # Detect horizontal pixel segments where characters are likely present.
        character_pixel_segments = get_projection_segments(roi, color, axis=0, min_length=2)
        if not character_pixel_segments:
            return []

        char_objects = []
        min_char_pixel_width = 4

        # Unpack the tight bounding box in JSON coordinates.
        json_x1, json_y1, _, json_y2 = tight_bbox

        # Iterate over each detected pixel segment.
        for start_px, end_px in character_pixel_segments:
            if end_px - start_px < min_char_pixel_width:
                continue

            # Convert the relative pixel coordinates of the segment back to absolute JSON coordinates.
            char_json_x1 = json_x1 + start_px * scale_x
            char_json_x2 = json_x1 + end_px * scale_x

            # Create the final bounding box for the character.
            char_bbox = [char_json_x1, json_y1, char_json_x2, json_y2]
            char_objects.append(Character(bbox=char_bbox, color=color, line_index=line_index))

        return char_objects

    def _analyze_and_correct_bboxes(self, char_objects, full_text, coords):
        non_space_chars = [c for c in full_text if c not in " \n"]
        expected_count = len(non_space_chars)
        if not char_objects or len(char_objects) < expected_count:
            return char_objects

        chars = sorted(char_objects, key=lambda c: (c.line_index, c.bbox[0]))
        num_fragments = len(chars)
        num_chars = len(non_space_chars)

        try:
            font = ImageFont.truetype("msyh.ttc", size=30)
            ideal_height = 30
            ideal_char_ratios = []
            for c in non_space_chars:
                if c in "iI,":
                    ideal_char_ratios.append(0.15)
                elif c in "，。、；：？！（）‘’":
                    ideal_char_ratios.append(0.25)
                elif c in "【】“”《》":
                    ideal_char_ratios.append(0.35)
                else:
                    ideal_char_ratios.append(font.getlength(c) / ideal_height)
        except IOError:
            return chars

        memo_cost = {}

        def get_merge_cost(start, end, char_idx):
            if (start, end, char_idx) in memo_cost:
                return memo_cost[(start, end, char_idx)]

            # Rule 2: Don't merge boxes of different colors.
            first_color = chars[start].color
            for i in range(start + 1, end):
                if chars[i].color != first_color:
                    memo_cost[(start, end, char_idx)] = float('inf')
                    return float('inf')

            # Rule 1: Don't merge boxes of different heights (with 20% tolerance).
            heights = [c.bbox[3] - c.bbox[1] for c in chars[start:end]]
            if not heights:
                memo_cost[(start, end, char_idx)] = float('inf')
                return float('inf')
            min_h, max_h = min(heights), max(heights)
            if max_h > min_h * 1.2:
                memo_cost[(start, end, char_idx)] = float('inf')
                return float('inf')

            merged_bbox = [chars[start].bbox[0],
                           min(c.bbox[1] for c in chars[start:end]),
                           chars[end - 1].bbox[2],
                           max(c.bbox[3] for c in chars[start:end])]

            merged_width = merged_bbox[2] - merged_bbox[0]
            merged_height = merged_bbox[3] - merged_bbox[1]
            if merged_height == 0:
                return float('inf')

            merged_ratio = merged_width / merged_height
            ideal_ratio = ideal_char_ratios[char_idx]

            # If ideal ratio is for a full-width char (close to 1) and the detected ratio
            # is slightly narrower (0.9-1.0), treat it as a perfect match (cost 0).
            if ideal_ratio > 0.9 and 0.9 <= merged_ratio <= 1.0:
                cost = 0
            else:
                cost = abs(merged_ratio - ideal_ratio)

            # Adjust cost based on the gap with the preceding character. A larger gap reduces the cost.
            gap_width = 0
            if start > 0:
                gap = chars[start].bbox[0] - chars[start - 1].bbox[2]
                if gap > 0:
                    gap_width = gap

            if merged_height > 0:
                cost -= 0.1 * (gap_width / merged_height)

            memo_cost[(start, end, char_idx)] = cost
            return cost

        dp = [[float('inf')] * (num_chars + 1) for _ in range(num_fragments + 1)]
        path = [[0] * (num_chars + 1) for _ in range(num_fragments + 1)]
        dp[0][0] = 0

        for j in range(1, num_chars + 1):
            for i in range(1, num_fragments + 1):
                for k in range(i):
                    cost = get_merge_cost(k, i, j - 1)
                    if dp[k][j - 1] + cost < dp[i][j]:
                        dp[i][j] = dp[k][j - 1] + cost
                        path[i][j] = k

        if dp[num_fragments][num_chars] == float('inf'):
            return []

        final_chars = []
        curr_frag = num_fragments
        for curr_char in range(num_chars, 0, -1):
            prev_frag = path[curr_frag][curr_char]

            merged_bbox = [chars[prev_frag].bbox[0],
                           min(c.bbox[1] for c in chars[prev_frag:curr_frag]),
                           chars[curr_frag - 1].bbox[2],
                           max(c.bbox[3] for c in chars[prev_frag:curr_frag])]

            new_char = Character(merged_bbox, chars[prev_frag].color, chars[prev_frag].line_index)
            new_char.text = non_space_chars[curr_char - 1]
            final_chars.append(new_char)
            curr_frag = prev_frag

        return final_chars[::-1]

    def _normalize_font_sizes(self, styles):
        if not styles:
            return styles

        i = 0
        while i < len(styles):
            j = i
            while j + 1 < len(styles) and abs(styles[j + 1].font_size - styles[j].font_size) < 3:
                j += 1

            group = styles[i:j + 1]
            if group:
                sizes = [s.font_size for s in group]
                most_common_size = Counter(sizes).most_common(1)[0][0]
                for style in group:
                    style.font_size = most_common_size

            i = j + 1

        return styles

    def _normalize_colors(self, styles, threshold=50):
        if not styles:
            return styles

        i = 0
        while i < len(styles):
            j = i
            while j + 1 < len(styles) and np.linalg.norm(
                    np.array(styles[j + 1].color) - np.array(styles[j].color)) < threshold:
                j += 1

            group = styles[i:j + 1]
            if group:
                colors = [tuple(s.color) for s in group]
                most_common_color = Counter(colors).most_common(1)[0][0]
                for style in group:
                    style.color = most_common_color
            i = j + 1
        return styles

    def _determine_character_styles(self, final_chars, coords, elem_type):
        for char in final_chars:
            height_pts = (char.bbox[3] - char.bbox[1]) * 0.95 * coords['scale_y']
            char.font_size = int(max(height_pts, 6.0))
            char.bold = elem_type == "title"
        return final_chars

    def _draw_debug_boxes_for_page(self, image, all_chars, coords, output_path):
        """Draws bounding boxes for an entire page's characters for debugging."""
        debug_img = image.copy()
        for char in all_chars:
            bbox = char.bbox
            px_box = [
                int(bbox[0] * coords['img_w'] / coords['json_w']),
                int(bbox[1] * coords['img_h'] / coords['json_h']),
                int(bbox[2] * coords['img_w'] / coords['json_w']),
                int(bbox[3] * coords['img_h'] / coords['json_h'])
            ]
            cv2.rectangle(debug_img, (px_box[0], px_box[1]), (px_box[2], px_box[3]), (0, 0, 255), 2)  # Red box
        cv2.imwrite(output_path, cv2.cvtColor(debug_img, cv2.COLOR_RGB2BGR))

    def _process_text(self, context, elem):
        bbox = elem.get("bbox")
        if not bbox: return

        context.add_element_bbox_for_cleanup(bbox)

        all_spans = [s for l in elem.get("lines", []) for s in l.get("spans", [])] if "lines" in elem else elem.get(
            "spans", [])
        if not all_spans:
            text_content = elem.get("text", "")
            if text_content:
                context.add_processed_element('text', {'bbox': bbox, 'text_runs': [{'text': text_content}]})
            return

        # Pre-process spans for bullet points
        if all_spans:
            first_span_content = all_spans[0].get("content", "")
            if first_span_content.lstrip().startswith('-'):
                cleaned_content = first_span_content.lstrip(' \t\n\r\f\v-•*·')
                all_spans[0]["content"] = "• " + cleaned_content

        full_text = "".join([s.get("content", "").replace('\\%', '%') for s in all_spans])
        if not full_text.strip(): return
        print(f"\n--- Processing Text ---\nContent: '{full_text.strip()[:100]}...'")

        try:
            line_infos = self._get_line_ranges(context.original_image, bbox, context.coords)
            print(f"Detected lines: {len(line_infos)}")
            if not line_infos: raise ValueError("No lines detected.")

            raw_chars = self._detect_raw_characters(context.original_image, line_infos, bbox, context.coords)

            # Heuristic for fixing missing bullet points based on color change
            if len(raw_chars) >= 2 and raw_chars[0].color != raw_chars[1].color:
                if not full_text.lstrip().startswith(('•', '·', '*', '-')):
                    full_text = "• " + full_text

            corrected_chars = self._analyze_and_correct_bboxes(raw_chars, full_text, context.coords)
            context.add_characters(raw_chars, corrected_chars)

            non_space_chars = [c for c in full_text if c not in " \n"]
            can_align = len(corrected_chars) == len(non_space_chars)
            print(f"Character alignment successful: {can_align}")
            print(f"Using char-by-char styling (mixed layout support): {can_align}")

            if not can_align:
                # Fallback to line-based rendering
                text_runs = self._get_text_runs_by_line(all_spans, line_infos, context.coords, elem.get("type"))
                context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs})
                return

            final_styles = self._determine_character_styles(corrected_chars, context.coords, elem.get("type"))
            final_styles = self._normalize_font_sizes(final_styles)
            final_styles = self._normalize_colors(final_styles)

            # Determine if single-line based on character analysis BEFORE adding to context
            line_indices = {char.line_index for char in final_styles}
            is_single_line = len(line_indices) <= 1

            text_runs = self._get_text_runs_by_char(full_text, final_styles)
            context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs, 'is_single_line': is_single_line})

        except Exception:
            # Broad exception fallback
            text_runs = self._get_text_runs_from_spans(all_spans, bbox, context.original_image, context.coords, elem.get("type"))
            # Fallback check for single-line
            is_single_line = not any('\n' in run['text'] for run in text_runs)
            context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs, 'is_single_line': is_single_line})

    def _get_text_runs_by_char(self, full_text, final_styles):
        """Generates styled text runs from character-by-character analysis."""
        runs = []
        style_iter = iter(final_styles)
        last_style = None
        for char in full_text:
            font_info = {'name': "Microsoft YaHei"}
            if char not in " \n":
                style = next(style_iter, None)
                if style:
                    font_info['size'] = Pt(style.font_size)
                    font_info['color'] = RGBColor(*style.color)
                    font_info['bold'] = style.bold
                    last_style = style
            elif last_style:
                font_info['size'] = Pt(last_style.font_size)
                font_info['color'] = RGBColor(*last_style.color)
                font_info['bold'] = last_style.bold
            runs.append({'text': char, 'font': font_info})
        return runs

    def _get_text_runs_by_line(self, all_spans, line_infos, coords, elem_type):
        """Generates styled text runs using line-based analysis as a fallback."""
        runs = []
        span_idx = 0
        for i, info in enumerate(line_infos):
            line_range = info['range']
            line_spans = []
            while span_idx < len(all_spans):
                span = all_spans[span_idx]
                sbbox = span.get("bbox")
                if sbbox and sbbox[1] < line_range[1] and sbbox[3] > line_range[0]:
                    line_spans.append(span)
                    span_idx += 1
                else:
                    break
            if not line_spans: continue

            line_text = "".join([s.get("content", "").replace('\\%', '%') for s in line_spans])
            if not line_text.strip() and i < len(line_infos) - 1: line_text += "\n"

            font_size_pts = (line_range[1] - line_range[0]) * coords['scale_y']
            font_info = {
                'name': "Microsoft YaHei",
                'color': RGBColor(*info['color']),
                'size': Pt(int(max(font_size_pts, 6.0))),
                'bold': elem_type == "title"
            }
            runs.append({'text': line_text, 'font': font_info})
            if i < len(line_infos) - 1 and not line_text.endswith('\n'):
                runs.append({'text': '\n', 'font': font_info}) # Keep consistent font for newline
        return runs

    def _get_text_runs_from_spans(self, spans, bbox, page_image, coords, elem_type=None):
        """Generates a single styled text run as a last-resort fallback."""
        if not spans: return []
        font_size_pts = (bbox[3] - bbox[1]) * coords['scale_y']
        full_text = "".join([s.get("content", "").replace('\\%', '%') for s in spans])
        bg_color = extract_background_color(page_image, bbox)
        color, _, _ = extract_font_color(page_image, bbox, bg_color)
        font_info = {
            'name': "Microsoft YaHei",
            'size': Pt(int(font_size_pts)),
            'bold': elem_type == "title",
            'color': RGBColor(*color)
        }
        return [{'text': full_text, 'font': font_info}]

    def _render_text_from_data(self, slide, text_data):
        """Renders a text element from processed data onto a slide."""
        bbox = text_data['bbox']
        text_runs = text_data.get('text_runs', [])
        is_single_line = text_data.get('is_single_line', False)

        # If it's a single line, widen the textbox to prevent wrapping due to font differences.
        if is_single_line:
            x1, y1, x2, y2 = bbox
            width = x2 - x1
            new_x2 = x1 + width * 1.2
            render_bbox = [x1, y1, new_x2, y2]
        else:
            render_bbox = bbox

        txBox = self._create_textbox(slide, render_bbox, self.coords_for_render)
        tf = txBox.text_frame
        tf.clear()
        tf.margin_bottom = tf.margin_top = tf.margin_left = tf.margin_right = Pt(0)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        for run_data in text_runs:
            run = p.add_run()
            run.text = run_data['text']
            font = run.font
            font_info = run_data.get('font', {})
            font.name = font_info.get('name', "Microsoft YaHei")
            if 'size' in font_info: font.size = font_info['size']
            if 'color' in font_info: font.color.rgb = font_info['color']
            if 'bold' in font_info: font.bold = font_info['bold']

    def _add_picture_from_bbox(self, slide, bbox, page_image, coords, text_elements):
        if not bbox: return
        x1, y1, x2, y2 = bbox;
        left, top, w, h = Pt(x1 * coords['scale_x']), Pt(y1 * coords['scale_y']), Pt((x2 - x1) * coords['scale_x']), Pt(
            (y2 - y1) * coords['scale_y'])
        px_box = [int(x1 * coords['img_w'] / coords['json_w']), int(y1 * coords['img_h'] / coords['json_h']),
                  int(x2 * coords['img_w'] / coords['json_w']), int(y2 * coords['img_h'] / coords['json_h'])]
        crop = page_image[px_box[1]:px_box[3], px_box[0]:px_box[2]].copy()

        # This cleanup logic is now less critical due to the global background inpainting,
        # but can still be useful for images that contain text not defined as a separate text element.
        for txt_e in text_elements:
            txt_box = txt_e.get("bbox")
            if txt_box and self._get_bbox_intersection(bbox, txt_box):
                px_txt_box = [int(v * (
                    coords['img_w'] / coords['json_w'] if i % 2 == 0 else coords['img_h'] / coords['json_h'])) for
                              i, v in enumerate(txt_box)]
                inter = self._get_bbox_intersection(px_box, px_txt_box)
                if inter:
                    local_inter = [inter[0] - px_box[0], inter[1] - px_box[1], inter[2] - px_box[0],
                                   inter[3] - px_box[1]]
                    fill_bbox_with_bg(crop, local_inter)

        if crop.size > 0:
            path = f"temp_crop_img_{x1}_{y1}.png";
            cv2.imwrite(path, cv2.cvtColor(crop, cv2.COLOR_RGB2BGR))
            slide.shapes.add_picture(path, left, top, w, h);
            os.remove(path)

    def _process_image(self, context, elem, text_elements):
        context.add_element_bbox_for_cleanup(elem.get("bbox"))

        if "blocks" in elem and elem["blocks"]:
            image_block_bbox = None
            for block in elem["blocks"]:
                if block.get("type") == "image_body" or (block.get("spans") and block["spans"][0].get("type") == "image"):
                    image_block_bbox = block.get("bbox")
                    break

            # Add the main image part to the render queue
            img_bbox_to_render = image_block_bbox or elem.get("bbox")
            context.add_processed_element('image', {'bbox': img_bbox_to_render, 'text_elements': text_elements})

            # Process any text blocks within the image element
            for block in elem["blocks"]:
                if block.get("type") == "image_caption":
                    self._process_text(context, block)
        else:
            # Simple image
            context.add_processed_element('image', {'bbox': elem.get("bbox"), 'text_elements': text_elements})

    def _process_list(self, context, elem):
        for block in elem.get("blocks", []):
            # Prepend bullet point
            spans = [s for l in block.get("lines", []) for s in l.get("spans", [])] if "lines" in block else block.get("spans", [])
            if spans:
                spans.sort(key=lambda s: (s.get("bbox", [0,0,0,0])[1], s.get("bbox", [0,0,0,0])[0]))
                spans[0]["content"] = "• " + spans[0].get("content", "").lstrip(' ·-*•')

            # Re-assign modified spans back to the block before processing
            if "lines" in block and block["lines"]:
                block["lines"][0]["spans"] = spans
            else:
                block["spans"] = spans

            self._process_text(context, block)

    def _process_element(self, context, elem, all_text_elements):
        cat = elem.get("type", "text")
        if cat == "list":
            self._process_list(context, elem)
        elif cat in ["text", "title", "caption", "footnote", "footer", "header", "page_number"]:
            self._process_text(context, elem)
        elif cat in ["image", "table", "formula", "figure"]:
            self._process_image(context, elem, all_text_elements)

    def process_page(self, slide, elements, page_image, page_size=None, page_index=0, debug_images=False):
        self.debug_images = debug_images
        img_h, img_w = page_image.shape[:2]
        json_w, json_h = page_size if page_size and all(page_size) else (img_w * 72 / 300, img_h * 72 / 300)
        w_pts, h_pts = self.cap_size(json_w, json_h)
        self.prs.slide_width, self.prs.slide_height = Pt(w_pts), Pt(h_pts)
        coords = {'scale_x': w_pts / json_w, 'scale_y': h_pts / json_h, 'img_w': img_w, 'img_h': img_h,
                  'json_w': json_w, 'json_h': json_h}
        self.coords_for_render = coords # Store for render phase

        context = PageContext(page_image, coords, slide)

        text_types = ["list", "text", "title", "caption", "footnote", "footer", "header", "page_number"]
        all_text_elements = [e for e in elements if e.get("type", "text") in text_types]

        # Phase 1: Analyze and populate context.
        # First, register bboxes for background cleanup.
        # Erase all elements, or just discarded blocks if watermark removal is on.
        for elem in elements:
            is_discarded = elem.get('is_discarded', False)
            if not is_discarded or (is_discarded and self.remove_watermark):
                context.add_element_bbox_for_cleanup(elem.get("bbox"))

        # Second, process elements to extract content for rendering.
        for elem in elements:
            # If watermark removal is on, skip processing/rendering discarded blocks.
            if elem.get('is_discarded') and self.remove_watermark:
                continue
            self._process_element(context, elem, all_text_elements)

        # Phase 2: Render from context
        context.render_to_slide(self)

        # Phase 3: Generate debug output if enabled
        if self.debug_images:
            context.generate_debug_images(page_index, self)

    def save(self):
        self.prs.save(self.output_path)


def convert_mineru_to_ppt(json_path, input_path, output_ppt_path, remove_watermark=True, debug_images=False):
    from .utils import pdf_to_images
    DPI = 300

    if debug_images:
        if os.path.exists("tmp"):
            shutil.rmtree("tmp")
        os.makedirs("tmp")

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Generate page images from either a PDF or a single image file
    if input_path.lower().endswith('.pdf'):
        images = pdf_to_images(input_path, dpi=DPI)
    else:
        try:
            img = Image.open(input_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            images = [np.array(img)]
        except Exception as e:
            raise IOError(f"Failed to load image file: {input_path} - {e}")

    gen = PPTGenerator(output_ppt_path, remove_watermark=remove_watermark)
    pages = data if isinstance(data, list) else next(
        (data[k] for k in ["pdf_info", "pages"] if k in data and isinstance(data[k], list)), [data])
    print(f"[CLEANUP] Found {len(pages)} pages.")
    for i, page_data in enumerate(pages):
        if i >= len(images): break
        print(f"Processing page {i + 1}/{len(pages)}...")
        page_img = images[i].copy()
        if i == 0: gen.set_slide_size(page_img.shape[1], page_img.shape[0], dpi=DPI)
        slide = gen.add_slide()

        elements = []
        for key in ["para_blocks", "images", "tables"]:
            for item in page_data.get(key, []):
                item['is_discarded'] = False
                elements.append(item)
        for item in page_data.get("discarded_blocks", []):
            item['is_discarded'] = True
            elements.append(item)

        page_size = page_data.get("page_size") or (page_data.get("page_info", {}).get("width"),
                                                   page_data.get("page_info", {}).get("height"))
        gen.process_page(slide, elements, page_img, page_size=page_size, page_index=i, debug_images=debug_images)
    gen.save()
    print(f"Saved to {output_ppt_path}")
